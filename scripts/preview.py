#!/usr/bin/env python3
"""sketch-to-slide preview: draws a rendered .pptx to a PNG for visual QA.

Approximate wireframe of what PowerPoint will show (shapes, text, connectors).
Requires Pillow. Elbow connectors are drawn with the same orthogonal routing
PowerPoint derives from the glued connection sites.
"""

import argparse
import sys

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.shapes.connector import Connector

W = 1280
FONT_PATHS = [
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    "/System/Library/Fonts/Helvetica.ttc",
]


def load_font(px):
    for p in FONT_PATHS:
        try:
            return ImageFont.truetype(p, px)
        except OSError:
            continue
    return ImageFont.load_default()


def draw_arrowhead(d, tip, src, size=10):
    import math
    ang = math.atan2(tip[1] - src[1], tip[0] - src[0])
    for da in (2.6, -2.6):
        d.line([tip, (tip[0] + size * math.cos(ang + da), tip[1] + size * math.sin(ang + da))],
               fill="black", width=2)


def elbow_points(p1, p2):
    dx, dy = abs(p2[0] - p1[0]), abs(p2[1] - p1[1])
    if dx >= dy:
        mx = (p1[0] + p2[0]) / 2
        return [p1, (mx, p1[1]), (mx, p2[1]), p2]
    my = (p1[1] + p2[1]) / 2
    return [p1, (p1[0], my), (p2[0], my), p2]


def draw_text_block(d, text, cx, cy, w, font, anchor_middle=True, align_center=True, top=None):
    lines = text.split("\n")
    lh = font.size + 4
    total = lh * len(lines)
    y = cy - total / 2 if anchor_middle else (top if top is not None else cy)
    for ln in lines:
        tw = d.textlength(ln, font=font)
        x = cx - tw / 2 if align_center else cx - w / 2 + 6
        d.text((x, y), ln, fill="black", font=font)
        y += lh


def render_preview(pptx_path, out_path):
    prs = Presentation(pptx_path)
    sw, sh = prs.slide_width, prs.slide_height
    H = int(W * sh / sw)
    img = Image.new("RGB", (W, H), "white")
    d = ImageDraw.Draw(img)

    def px(v):
        return v / sw * W

    def py(v):
        return v / sh * H

    slide = prs.slides[0]
    conns, shapes = [], []
    for sp in slide.shapes:
        (conns if isinstance(sp, Connector) else shapes).append(sp)

    for sp in conns:
        p1 = (px(sp.begin_x), py(sp.begin_y))
        p2 = (px(sp.end_x), py(sp.end_y))
        prst = sp._element.find(qn("p:spPr")).find(qn("a:prstGeom"))
        bent = prst is not None and "bent" in (prst.get("prst") or "")
        pts = elbow_points(p1, p2) if bent else [p1, p2]
        d.line(pts, fill="black", width=2)
        ln = sp.line._get_or_add_ln()
        if ln.find(qn("a:tailEnd")) is not None:
            draw_arrowhead(d, pts[-1], pts[-2])
        if ln.find(qn("a:headEnd")) is not None:
            draw_arrowhead(d, pts[0], pts[1])

    for sp in shapes:
        x, y = px(sp.left), py(sp.top)
        w, h = px(sp.width), py(sp.height)
        try:
            auto = sp.auto_shape_type
        except Exception:
            auto = None
        rot = getattr(sp, "rotation", 0) or 0
        box = (x, y, x + w, y + h)
        if rot in (90, 270) and auto in (MSO_SHAPE.CHEVRON, MSO_SHAPE.ISOSCELES_TRIANGLE):
            cx, cy = x + w / 2, y + h / 2
            w, h = h, w
            x, y = cx - w / 2, cy - h / 2
            box = (x, y, x + w, y + h)
        if auto == MSO_SHAPE.OVAL:
            d.ellipse(box, outline="black", width=2)
        elif auto == MSO_SHAPE.DIAMOND:
            d.polygon([(x + w / 2, y), (x + w, y + h / 2), (x + w / 2, y + h), (x, y + h / 2)],
                      outline="black", width=2)
        elif auto == MSO_SHAPE.ISOSCELES_TRIANGLE:
            if rot == 180:
                d.polygon([(x, y), (x + w, y), (x + w / 2, y + h)], outline="black", width=2)
            elif rot == 90:
                d.polygon([(x, y), (x + w, y + h / 2), (x, y + h)], outline="black", width=2)
            elif rot == 270:
                d.polygon([(x + w, y), (x + w, y + h), (x, y + h / 2)], outline="black", width=2)
            else:
                d.polygon([(x + w / 2, y), (x + w, y + h), (x, y + h)], outline="black", width=2)
        elif auto == MSO_SHAPE.CHEVRON:
            n = w * 0.35
            if rot == 90:
                n = h * 0.35
                d.polygon([(x, y), (x + w / 2, y + n), (x + w, y), (x + w, y + h - n),
                           (x + w / 2, y + h), (x, y + h - n)], outline="black", width=2)
            else:
                d.polygon([(x, y), (x + w - n, y), (x + w, y + h / 2), (x + w - n, y + h),
                           (x, y + h), (x + n, y + h / 2)], outline="black", width=2)
        elif auto == MSO_SHAPE.ROUNDED_RECTANGLE:
            d.rounded_rectangle(box, radius=min(w, h) * 0.15, outline="black", width=2)
        elif auto == MSO_SHAPE.RECTANGLE:
            d.rectangle(box, outline="black", width=2)

        if sp.has_text_frame and sp.text_frame.text.strip():
            tf = sp.text_frame
            first = tf.paragraphs[0]
            size_pt = None
            bold = False
            for p in tf.paragraphs:
                for r in p.runs:
                    if size_pt is None and r.font.size is not None:
                        size_pt = r.font.size.pt
                        bold = bool(r.font.bold)
                if size_pt is not None:
                    break
            size_pt = size_pt or 12
            fpx = max(9, int(size_pt / 7.5 * H / 100 * 1.55))
            font = load_font(fpx)
            from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
            middle = tf.vertical_anchor in (MSO_ANCHOR.MIDDLE, None)
            centered = first.alignment in (PP_ALIGN.CENTER, None) and auto is not None
            if first.alignment == PP_ALIGN.LEFT:
                centered = False
            txt = tf.text
            if bold:
                d.text((x + (w - d.textlength(txt.split("\n")[0], font=font)) / 2 if centered else x + 4,
                        y + (h - font.size) / 2 if middle else y + 4),
                       txt, fill="black", font=font, stroke_width=1, stroke_fill="black")
            else:
                draw_text_block(d, txt, x + w / 2, y + h / 2, w, font,
                                anchor_middle=middle, align_center=centered, top=y + 4)

    img.save(out_path)
    print(out_path)


def main():
    ap = argparse.ArgumentParser(description="Draw an approximate PNG preview of a rendered .pptx")
    ap.add_argument("pptx")
    ap.add_argument("-o", "--output")
    args = ap.parse_args()
    out = args.output or (args.pptx.rsplit(".", 1)[0] + "-preview.png")
    render_preview(args.pptx, out)


if __name__ == "__main__":
    main()

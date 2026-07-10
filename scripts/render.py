#!/usr/bin/env python3
"""sketch-to-slide renderer: layout JSON -> 16:9 .pptx of native PowerPoint shapes.

Applies a deterministic tidy pass (canvas fit, row/column alignment, size
unification) so hand-estimated coordinates come out clean without changing
the sketch's structure.
"""

import argparse
import copy
import json
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

SLIDE_W_EMU = 12192000
SLIDE_H_EMU = 6858000
ASPECT = (SLIDE_W_EMU / SLIDE_H_EMU)  # 1.777...; converts normalized dx to visual units
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
SIZE_PT = {"title": 20, "body": 12, "small": 9}
ALIGN_MAP = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
SHAPE_MAP = {
    "rect": MSO_SHAPE.RECTANGLE,
    "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_SHAPE.OVAL,
    "circle": MSO_SHAPE.OVAL,
    "diamond": MSO_SHAPE.DIAMOND,
    "chevron_right": MSO_SHAPE.CHEVRON,
    "chevron_down": MSO_SHAPE.CHEVRON,
    "triangle_up": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "triangle_right": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "triangle_down": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "triangle_left": MSO_SHAPE.ISOSCELES_TRIANGLE,
}
ROTATION = {"chevron_down": 90, "triangle_right": 90, "triangle_down": 180, "triangle_left": 270}
CONNECTOR_TYPES = ("line", "arrow", "double_arrow")
TITLE_BOX = (0.04, 0.035, 0.92, 0.10)

MARGIN_X = (0.05, 0.95)
MARGIN_Y_TITLED = (0.17, 0.93)
MARGIN_Y_PLAIN = (0.06, 0.94)
MAX_UPSCALE = 2.2
ROW_TOL = 0.05
COL_TOL = 0.035
SIZE_TOL_W = 0.035
SIZE_TOL_H = 0.05
STRAIGHT_TOL = 0.02
FLUSH_TOL = 0.015
GAP_RATIO_MAX = 2.0
DECOR_FRAC = 0.5
HEADER_H = 0.055
FRAME_H_FRAC = 0.55
FRAME_COVER_FRAC = 0.70
RATIO_TOL = 0.08
FRAME_CANONS = ([1 / 3, 2 / 3], [2 / 3, 1 / 3], [1 / 3, 1 / 3, 1 / 3], [0.25, 0.5, 0.25])


def fail(msg):
    print(json.dumps({"error": msg}), file=sys.stderr)
    sys.exit(2)


def load_layout(path):
    try:
        with open(path) as f:
            layout = json.load(f)
    except (OSError, json.JSONDecodeError) as e:
        fail(f"cannot read layout: {e}")
    if layout.get("version") != 1:
        fail("layout version must be 1")
    elements = layout.get("elements")
    if not isinstance(elements, list) or not elements:
        fail("layout needs a non-empty 'elements' list")
    seen = set()
    table_ids = set()
    for i, el in enumerate(elements):
        t = el.get("type")
        if t not in SHAPE_MAP and t not in ("text", "table") and t not in CONNECTOR_TYPES:
            fail(f"element {i}: unknown type '{t}'")
        eid = el.get("id")
        if eid is not None:
            if eid in seen:
                fail(f"duplicate id '{eid}'")
            seen.add(eid)
        if t == "circle":
            for k in ("x", "y", "w"):
                if not isinstance(el.get(k), (int, float)):
                    fail(f"element {i} ('circle'): missing numeric '{k}'")
        elif t in SHAPE_MAP or t in ("text", "table"):
            for k in ("x", "y", "w", "h"):
                if not isinstance(el.get(k), (int, float)):
                    fail(f"element {i} ('{t}'): missing numeric '{k}'")
        elif "from" not in el and "to" not in el:
            for k in ("x1", "y1", "x2", "y2"):
                if not isinstance(el.get(k), (int, float)):
                    fail(f"element {i} ('{t}'): needs from/to ids or x1,y1,x2,y2")
        if t == "table":
            rows = el.get("rows")
            if (not isinstance(rows, list) or not rows
                    or not all(isinstance(r, list) and r and all(isinstance(c, str) for c in r) for r in rows)):
                fail(f"element {i} ('table'): 'rows' must be a non-empty list of lists of strings")
            cols = el.get("columns")
            if cols is not None and (not isinstance(cols, list) or not cols
                                     or not all(isinstance(c, str) for c in cols)):
                fail(f"element {i} ('table'): 'columns' must be a non-empty list of strings")
            ncols = len(cols) if cols else len(rows[0])
            if any(len(r) != ncols for r in rows):
                fail(f"element {i} ('table'): every row must have {ncols} cells")
            if eid is not None:
                table_ids.add(eid)
        b = el.get("bullets")
        if b is not None and (not isinstance(b, list) or not all(isinstance(s, str) for s in b)):
            fail(f"element {i}: 'bullets' must be a list of strings")
    for el in elements:
        if el.get("type") in CONNECTOR_TYPES:
            for key in ("from", "to"):
                if el.get(key) in table_ids:
                    fail(f"connector '{el.get('id', '?')}' references table '{el[key]}'; "
                         "connectors cannot attach to tables")
    return layout


def _is_box(el):
    return el["type"] in SHAPE_MAP or el["type"] in ("text", "table")


def _fix_circles(elements):
    for el in elements:
        if el["type"] == "circle":
            h = el["w"] * ASPECT
            cy = el["y"] + el.get("h", h) / 2
            el["h"] = h
            el["y"] = cy - h / 2


def _fit_canvas(layout):
    els = layout["elements"]
    pts = []
    for el in els:
        if _is_box(el):
            pts += [(el["x"], el["y"]), (el["x"] + el["w"], el["y"] + el["h"])]
        elif "x1" in el:
            pts += [(el["x1"], el["y1"]), (el["x2"], el["y2"])]
    if not pts:
        return 1.0, 1.0
    x0 = min(p[0] for p in pts)
    x1 = max(p[0] for p in pts)
    y0 = min(p[1] for p in pts)
    y1 = max(p[1] for p in pts)
    bw, bh = x1 - x0, y1 - y0
    mx0, mx1 = MARGIN_X
    my0, my1 = MARGIN_Y_TITLED if layout.get("title") else MARGIN_Y_PLAIN
    fx = min((mx1 - mx0) / bw, MAX_UPSCALE) if bw > 0.02 else 1.0
    fy = min((my1 - my0) / bh, MAX_UPSCALE) if bh > 0.02 else 1.0
    xoff = mx0 + ((mx1 - mx0) - bw * fx) / 2
    yoff = my0 + ((my1 - my0) - bh * fy) / 2

    def tx(v):
        return xoff + (v - x0) * fx

    def ty(v):
        return yoff + (v - y0) * fy

    for el in els:
        if _is_box(el):
            el["x"], el["y"] = tx(el["x"]), ty(el["y"])
            el["w"], el["h"] = el["w"] * fx, el["h"] * fy
        elif "x1" in el:
            el["x1"], el["y1"] = tx(el["x1"]), ty(el["y1"])
            el["x2"], el["y2"] = tx(el["x2"]), ty(el["y2"])
    return fx, fy


def _cluster(values, tol):
    """Greedy 1-D clustering; returns list of (mean, indices)."""
    order = sorted(range(len(values)), key=lambda i: values[i])
    groups = []
    for i in order:
        if groups and abs(values[i] - groups[-1][0]) <= tol:
            mean, idxs = groups[-1]
            idxs.append(i)
            groups[-1] = (sum(values[j] for j in idxs) / len(idxs), idxs)
        else:
            groups.append((values[i], [i]))
    return groups


def _align(boxes):
    if len(boxes) < 2:
        return
    cys = [b["y"] + b["h"] / 2 for b in boxes]
    for mean, idxs in _cluster(cys, ROW_TOL):
        if len(idxs) > 1:
            for i in idxs:
                boxes[i]["y"] = mean - boxes[i]["h"] / 2
    cxs = [b["x"] + b["w"] / 2 for b in boxes]
    for mean, idxs in _cluster(cxs, COL_TOL):
        if len(idxs) > 1:
            for i in idxs:
                boxes[i]["x"] = mean - boxes[i]["w"] / 2


def _normalize_headers(boxes):
    """header: true boxes become slim section-header bands of fixed height,
    anchored at their drawn bottom edge so a drawn-flush seat on the body box
    survives; flush snap closes any remaining drawn gap.
    """
    for b in boxes:
        if b.get("header") and b["type"] in SHAPE_MAP:
            b["y"] = b["y"] + b["h"] - HEADER_H
            b["h"] = HEADER_H


def _unify_sizes(boxes):
    shapes = [b for b in boxes if b["type"] in SHAPE_MAP]
    if len(shapes) < 2:
        return
    # widths include headers so a band keeps tracking its frame's width;
    # heights exclude them — band height is contract-fixed and must never be
    # pooled with content boxes (field test 2026-07-10: chip pooled with squares)
    ws = [b["w"] for b in shapes]
    for mean, idxs in _cluster(ws, SIZE_TOL_W):
        if len(idxs) > 1:
            for i in idxs:
                b = shapes[i]
                cx = b["x"] + b["w"] / 2
                b["w"] = mean
                b["x"] = cx - mean / 2
    shapes = [b for b in shapes if not b.get("header")]
    if len(shapes) < 2:
        return
    hs = [b["h"] for b in shapes]
    for mean, idxs in _cluster(hs, SIZE_TOL_H):
        if len(idxs) > 1:
            for i in idxs:
                b = shapes[i]
                cy = b["y"] + b["h"] / 2
                b["h"] = mean
                b["y"] = cy - mean / 2


def _overlap(a0, a1, b0, b1):
    return max(0.0, min(a1, b1) - max(a0, b0))


def _snap_flush(boxes, fx, fy):
    """Boxes drawn touching or nearly touching snap flush (e.g. header chips).

    The contract's 1.5% bound is on the DRAWN gap, so the tolerance is scaled
    by the canvas-fit factors (fx, fy) — otherwise an upscaled small sketch
    misses snaps and a downscaled one snaps boxes that weren't near-touching.
    The bound is symmetric: slightly apart OR slightly overlapping both mean
    "drawn touching" (earlier passes like size unification can push a
    drawn-flush chip edge into its body box). The smaller box moves to close
    the gap; requires >=50% overlap on the perpendicular axis so unrelated
    neighbors are never pulled together.
    """
    for a in boxes:
        for b in boxes:
            if a is b:
                continue
            ov = _overlap(a["x"], a["x"] + a["w"], b["x"], b["x"] + b["w"])
            gap = b["y"] - (a["y"] + a["h"])
            # sum-of-heights guard: two boxes both thinner than the bound could
            # otherwise re-fire on the reverse pair and invert their stacking
            if (gap != 0 and abs(gap) <= FLUSH_TOL * fy and a["h"] + b["h"] > FLUSH_TOL * fy
                    and ov >= 0.5 * min(a["w"], b["w"])):
                small = a if a["w"] * a["h"] <= b["w"] * b["h"] else b
                small["y"] += gap if small is a else -gap
            ov = _overlap(a["y"], a["y"] + a["h"], b["y"], b["y"] + b["h"])
            gap = b["x"] - (a["x"] + a["w"])
            if (gap != 0 and abs(gap) <= FLUSH_TOL * fx and a["w"] + b["w"] > FLUSH_TOL * fx
                    and ov >= 0.5 * min(a["h"], b["h"])):
                small = a if a["w"] * a["h"] <= b["w"] * b["h"] else b
                small["x"] += gap if small is a else -gap


def _groups_by_center(boxes, axis):
    """Cluster boxes into column-groups (axis 0) or row-groups (axis 1).

    A group is the set of boxes sharing a center on that axis (a header chip
    and its body box form one column-group), so distribution moves them as a
    unit and never breaks their internal alignment.
    """
    p, s = ("x", "w") if axis == 0 else ("y", "h")
    centers = [b[p] + b[s] / 2 for b in boxes]
    tol = COL_TOL if axis == 0 else ROW_TOL
    groups = []
    for _mean, idxs in _cluster(centers, tol):
        grp = [boxes[i] for i in idxs]
        groups.append({
            "boxes": grp,
            "lo": min(b[p] for b in grp),
            "hi": max(b[p] + b[s] for b in grp),
        })
    groups.sort(key=lambda g: g["lo"])
    return groups


def _distribute(boxes, connectors, axis):
    """Equidistant distribution of 3+ aligned groups (contract: consulting
    exception 1). Bounded: gaps must already be similar (largest <= 2x
    smallest); a gap occupied by a small drawn shape or bridged by a
    connector is exempt from the bound. Uneven beyond the bound is
    deliberate and left untouched.
    """
    p = "x" if axis == 0 else "y"
    groups = _groups_by_center(boxes, axis)
    if len(groups) < 3:
        return
    extents = sorted(g["hi"] - g["lo"] for g in groups)
    median = extents[len(extents) // 2]
    members = [g for g in groups if g["hi"] - g["lo"] >= DECOR_FRAC * median]
    decor = [g for g in groups if g["hi"] - g["lo"] < DECOR_FRAC * median]
    if len(members) < 3:
        return
    for a, b in zip(members, members[1:]):
        if b["lo"] - a["hi"] < -0.005:  # ponytail: overlapping bands = not a clean banded layout; bail
            return
    member_ids = [
        {b.get("id") for b in g["boxes"] if b.get("id")} for g in members
    ]

    def bridged(i):
        return any(
            (c.get("from") in member_ids[i] and c.get("to") in member_ids[i + 1])
            or (c.get("to") in member_ids[i] and c.get("from") in member_ids[i + 1])
            for c in connectors
        )

    gaps = [b["lo"] - a["hi"] for a, b in zip(members, members[1:])]
    decor_gap = {}
    for d in decor:
        c = (d["lo"] + d["hi"]) / 2
        for i in range(len(members) - 1):
            if members[i]["hi"] <= c <= members[i + 1]["lo"]:
                decor_gap[id(d)] = i
                break
    free = [
        g for i, g in enumerate(gaps)
        if i not in decor_gap.values() and not bridged(i)
    ]
    if len(free) >= 2 and max(free) > GAP_RATIO_MAX * max(min(free), 1e-9):
        return
    span_lo, span_hi = members[0]["lo"], members[-1]["hi"]
    total = sum(g["hi"] - g["lo"] for g in members)
    gap = (span_hi - span_lo - total) / (len(members) - 1)
    if gap < 0:
        return
    pos = span_lo
    for g in members:
        dv = pos - g["lo"]
        for b in g["boxes"]:
            b[p] += dv
        size = g["hi"] - g["lo"]
        g["lo"], g["hi"] = pos, pos + size
        pos += size + gap
    for d in decor:
        i = decor_gap.get(id(d))
        if i is None:
            continue
        mid = (members[i]["hi"] + members[i + 1]["lo"]) / 2
        dv = mid - (d["lo"] + d["hi"]) / 2
        for b in d["boxes"]:
            b[p] += dv


CELL_PAD = 0.005


def _table_elements(el, idx):
    """Expand a table into native primitives: bold header boxes, borderless
    left-aligned text cells, and a horizontal separator line under each body
    row. Never a native PPT table object (contract rule).
    """
    x, y, w, h = el["x"], el["y"], el["w"], el["h"]
    cols = el.get("columns") or []
    rows = el["rows"]
    ncols = len(cols) if cols else len(rows[0])
    header = 1 if cols else 0
    rh = h / (len(rows) + header)
    cw = w / ncols
    eid = el.get("id") or f"table{idx}"
    out = []
    for j, head in enumerate(cols):
        out.append({"id": f"{eid}.h{j}", "type": "rect", "x": x + j * cw, "y": y,
                    "w": cw, "h": rh, "text": head, "header": True, "size": "small"})
    for i, row in enumerate(rows):
        ry = y + (header + i) * rh
        for j, cell in enumerate(row):
            if el.get("row_headers") and j == 0:
                out.append({"id": f"{eid}.r{i}c0", "type": "rect", "x": x, "y": ry,
                            "w": cw, "h": rh, "text": cell, "header": True, "size": "small"})
            else:
                out.append({"id": f"{eid}.r{i}c{j}", "type": "text",
                            "x": x + j * cw + CELL_PAD, "y": ry + CELL_PAD,
                            "w": cw - 2 * CELL_PAD, "h": rh - 2 * CELL_PAD,
                            "text": cell, "size": "small"})
        out.append({"id": f"{eid}.ln{i}", "type": "line",
                    "x1": x, "y1": ry + rh, "x2": x + w, "y2": ry + rh})
    return out


def _expand_tables(layout):
    out = []
    for i, el in enumerate(layout["elements"]):
        out.extend(_table_elements(el, i) if el["type"] == "table" else [el])
    layout["elements"] = out


def _snap_frames(boxes):
    """Canonical frame ratios (contract: consulting exception 3): 2-3 tall
    column-groups jointly spanning the content area, whose width fractions
    are within RATIO_TOL of a canonical consulting split, snap to the exact
    ratio. Widths scale per group; original gaps are preserved here and
    equalized afterwards by _distribute.

    ponytail: horizontal (left-to-right) frames only; vertical splits when a
    real sketch needs them.
    """
    if not boxes:
        return
    x0 = min(b["x"] for b in boxes)
    x1 = max(b["x"] + b["w"] for b in boxes)
    y0 = min(b["y"] for b in boxes)
    y1 = max(b["y"] + b["h"] for b in boxes)
    content_w, content_h = x1 - x0, y1 - y0
    if content_w <= 0 or content_h <= 0:
        return
    frames = []
    for g in _groups_by_center(boxes, 0):
        gy0 = min(b["y"] for b in g["boxes"])
        gy1 = max(b["y"] + b["h"] for b in g["boxes"])
        if gy1 - gy0 >= FRAME_H_FRAC * content_h:
            frames.append(g)
    if len(frames) not in (2, 3):
        return
    frames.sort(key=lambda g: g["lo"])
    for a, b in zip(frames, frames[1:]):
        if b["lo"] - a["hi"] < -0.005:
            return
    widths = [g["hi"] - g["lo"] for g in frames]
    total = sum(widths)
    if total < FRAME_COVER_FRAC * content_w:
        return
    fracs = [w / total for w in widths]
    best = None
    for canon in FRAME_CANONS:
        if len(canon) != len(frames):
            continue
        d = max(abs(f - c) for f, c in zip(fracs, canon))
        if d <= RATIO_TOL and (best is None or d < best[0]):
            best = (d, canon)
    if best is None:
        return
    gaps = [b["lo"] - a["hi"] for a, b in zip(frames, frames[1:])]
    pos = frames[0]["lo"]
    for i, (g, c) in enumerate(zip(frames, best[1])):
        new_w = c * total
        f = new_w / (g["hi"] - g["lo"])
        for b in g["boxes"]:
            b["x"] = pos + (b["x"] - g["lo"]) * f
            b["w"] *= f
        g["lo"], g["hi"] = pos, pos + new_w
        if i < len(gaps):
            pos += new_w + gaps[i]


def normalize_layout(layout):
    """Deterministic tidy pass shared by render and verify. Returns a deep copy."""
    layout = copy.deepcopy(layout)
    els = layout["elements"]
    _fix_circles(els)
    fx, fy = _fit_canvas(layout)
    boxes = [el for el in els if _is_box(el)]
    connectors = [el for el in els if el["type"] in CONNECTOR_TYPES]
    _normalize_headers(boxes)
    _align(boxes)
    _unify_sizes(boxes)
    _snap_flush(boxes, fx, fy)
    _snap_frames(boxes)
    _distribute(boxes, connectors, 0)
    _distribute(boxes, connectors, 1)
    _fix_circles(els)
    _expand_tables(layout)
    return layout


def boxes_by_id(layout):
    return {
        el["id"]: (el["x"], el["y"], el["w"], el["h"])
        for el in layout["elements"]
        if el.get("id") and el["type"] not in CONNECTOR_TYPES
    }


GLUE_TYPES = ("rect", "rounded_rect", "diamond")
SIDE_TOP, SIDE_LEFT, SIDE_BOTTOM, SIDE_RIGHT = 0, 1, 2, 3


def resolve_connector(el, boxes):
    """Return ((x1, y1), (x2, y2), kind, side1, side2) in normalized coordinates.

    Endpoints attach at side midpoints; kind is 'straight' when the boxes are
    aligned on the perpendicular axis, else 'elbow'. side1/side2 are PowerPoint
    connection-site indices (0=top, 1=left, 2=bottom, 3=right) or None for raw
    coordinate connectors.
    """
    if "from" not in el and "to" not in el:
        return (el["x1"], el["y1"]), (el["x2"], el["y2"]), "straight", None, None
    for key in ("from", "to"):
        if el.get(key) not in boxes:
            fail(f"connector '{el.get('id', '?')}' references unknown id '{el.get(key)}'")
    x1, y1, w1, h1 = boxes[el["from"]]
    x2, y2, w2, h2 = boxes[el["to"]]
    c1 = (x1 + w1 / 2, y1 + h1 / 2)
    c2 = (x2 + w2 / 2, y2 + h2 / 2)
    dx = (c2[0] - c1[0]) * ASPECT
    dy = c2[1] - c1[1]
    # overlap-aware side pick: vertically disjoint boxes (a row change) attach
    # bottom→top so the elbow routes through the inter-row gap instead of
    # cutting across boxes; same-row boxes attach left/right; only overlapping
    # boxes fall back to direction dominance
    vert_disjoint = (y1 + h1 <= y2) or (y2 + h2 <= y1)
    horiz_disjoint = (x1 + w1 <= x2) or (x2 + w2 <= x1)
    if vert_disjoint:
        horizontal = False
    elif horiz_disjoint:
        horizontal = True
    else:
        horizontal = abs(dx) >= abs(dy)
    if horizontal:
        p1 = (x1 + w1 if dx >= 0 else x1, c1[1])
        p2 = (x2 if dx >= 0 else x2 + w2, c2[1])
        s1 = SIDE_RIGHT if dx >= 0 else SIDE_LEFT
        s2 = SIDE_LEFT if dx >= 0 else SIDE_RIGHT
        kind = "straight" if abs(c1[1] - c2[1]) <= STRAIGHT_TOL else "elbow"
    else:
        p1 = (c1[0], y1 + h1 if dy >= 0 else y1)
        p2 = (c2[0], y2 if dy >= 0 else y2 + h2)
        s1 = SIDE_BOTTOM if dy >= 0 else SIDE_TOP
        s2 = SIDE_TOP if dy >= 0 else SIDE_BOTTOM
        kind = "straight" if abs(c1[0] - c2[0]) <= STRAIGHT_TOL else "elbow"
    return p1, p2, kind, s1, s2


def label_box(el, p1, p2):
    text = str(el["text"])
    w = min(0.30, max(0.10, 0.008 * len(text)))
    h = 0.05
    mx, my = (p1[0] + p2[0]) / 2, (p1[1] + p2[1]) / 2
    return mx - w / 2, my - h / 2, w, h


def _ex(v):
    return Emu(int(round(v * SLIDE_W_EMU)))


def _ey(v):
    return Emu(int(round(v * SLIDE_H_EMU)))


def expected_text(el):
    """Full text content of a rendered element (text plus bullet lines)."""
    lines = []
    if el.get("text"):
        lines.append(str(el["text"]))
    for b in el.get("bullets") or []:
        lines.append(f"• {b}")
    return "\n".join(lines)


def set_text(shape, el):
    tf = shape.text_frame
    tf.word_wrap = True
    bullets = el.get("bullets") or []
    text = el.get("text", "")
    size = el.get("size", "body")
    header = bool(el.get("header"))
    bold = True if header else el.get("bold", False)
    default_align = "left" if (header or el.get("type") == "text" or bullets) else "center"
    align = el.get("align", default_align)
    tf.vertical_anchor = MSO_ANCHOR.TOP if (el.get("type") == "text" or bullets) else MSO_ANCHOR.MIDDLE
    color = WHITE if header else BLACK
    pt = Pt(SIZE_PT.get(size, SIZE_PT["body"]))
    lines = []
    if text:
        lines += [(ln, bold) for ln in str(text).split("\n")]
    lines += [(f"• {b}", False) for b in bullets]
    if not lines:
        return
    for i, (line, is_bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = ALIGN_MAP.get(align, PP_ALIGN.CENTER)
        p.font.size = pt
        p.font.bold = bool(is_bold)
        p.font.color.rgb = color
        if bullets and i > 0:
            p.space_before = Pt(4)
        for r in p.runs:
            r.font.size = pt
            r.font.bold = bool(is_bold)
            r.font.color.rgb = color


def set_plain_text(shape, text, size="body", bold=False, align="center"):
    set_text(shape, {"text": text, "size": size, "bold": bold, "align": align})
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def _no_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def add_arrowheads(conn, head, tail):
    # python-pptx has no arrowhead API; append DrawingML line-end elements directly
    ln = conn.line._get_or_add_ln()
    if head:
        ln.append(ln.makeelement(qn("a:headEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    if tail:
        ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))


def _add_shape(slide, el):
    t = el["type"]
    rot = ROTATION.get(t, 0)
    x, y, w, h = el["x"], el["y"], el["w"], el["h"]
    ew, eh = _ex(w), _ey(h)
    left, top = _ex(x), _ey(y)
    if rot in (90, 270):
        cx, cy = left + ew / 2, top + eh / 2
        ew, eh = eh, ew
        left, top = int(cx - ew / 2), int(cy - eh / 2)
    sp = slide.shapes.add_shape(SHAPE_MAP[t], left, top, ew, eh)
    if rot:
        sp.rotation = rot
    return sp


def render(layout, out_path):
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W_EMU)
    prs.slide_height = Emu(SLIDE_H_EMU)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    layout = normalize_layout(layout)
    boxes = boxes_by_id(layout)
    counts = {"shapes": 0, "textboxes": 0, "connectors": 0, "labels": 0}

    if layout.get("title"):
        x, y, w, h = TITLE_BOX
        tb = slide.shapes.add_textbox(_ex(x), _ey(y), _ex(w), _ey(h))
        set_plain_text(tb, layout["title"], size="title", bold=True, align="left")
        counts["textboxes"] += 1

    connectors = []
    shapes_by_id = {}
    types_by_id = {}
    for el in layout["elements"]:
        t = el["type"]
        if t in CONNECTOR_TYPES:
            connectors.append(el)
            continue
        if t == "text":
            tb = slide.shapes.add_textbox(_ex(el["x"]), _ey(el["y"]), _ex(el["w"]), _ey(el["h"]))
            set_text(tb, el)
            if el.get("id"):
                shapes_by_id[el["id"]] = tb
                types_by_id[el["id"]] = t
            counts["textboxes"] += 1
            continue
        sp = _add_shape(slide, el)
        sp.fill.solid()
        sp.fill.fore_color.rgb = BLACK if el.get("header") else WHITE
        sp.line.color.rgb = BLACK
        sp.line.width = Pt(1)
        _no_shadow(sp)
        if el.get("text") or el.get("bullets"):
            set_text(sp, el)
        if el.get("id"):
            shapes_by_id[el["id"]] = sp
            types_by_id[el["id"]] = t
        counts["shapes"] += 1

    for el in connectors:
        p1, p2, kind, s1, s2 = resolve_connector(el, boxes)
        glue1 = types_by_id.get(el.get("from")) in GLUE_TYPES
        glue2 = types_by_id.get(el.get("to")) in GLUE_TYPES
        # bentConnector3 routes horizontal-first; vertical-dominant elbows only
        # look right when glued sites let PowerPoint re-route them
        horiz = s1 in (SIDE_LEFT, SIDE_RIGHT)
        use_elbow = kind == "elbow" and (horiz or (glue1 and glue2))
        ctype = MSO_CONNECTOR.ELBOW if use_elbow else MSO_CONNECTOR.STRAIGHT
        conn = slide.shapes.add_connector(ctype, _ex(p1[0]), _ey(p1[1]), _ex(p2[0]), _ey(p2[1]))
        if glue1:
            conn.begin_connect(shapes_by_id[el["from"]], s1)
        if glue2:
            conn.end_connect(shapes_by_id[el["to"]], s2)
        conn.line.color.rgb = BLACK
        conn.line.width = Pt(1.5)
        _no_shadow(conn)
        t = el["type"]
        add_arrowheads(conn, head=(t == "double_arrow"), tail=(t in ("arrow", "double_arrow")))
        counts["connectors"] += 1
        if el.get("text"):
            lx, ly, lw, lh = label_box(el, p1, p2)
            tb = slide.shapes.add_textbox(_ex(lx), _ey(ly), _ex(lw), _ey(lh))
            tb.fill.solid()
            tb.fill.fore_color.rgb = WHITE
            tb.line.fill.background()
            set_plain_text(tb, el["text"], "small", False, "center")
            counts["labels"] += 1

    prs.save(out_path)
    return counts


def main():
    ap = argparse.ArgumentParser(description="Render layout JSON into a 16:9 .pptx of native shapes")
    ap.add_argument("layout")
    ap.add_argument("-o", "--output")
    args = ap.parse_args()
    out = args.output or (args.layout.rsplit(".", 1)[0] + ".pptx")
    layout = load_layout(args.layout)
    counts = render(layout, out)
    print(json.dumps({"output": out, **counts}))


if __name__ == "__main__":
    main()

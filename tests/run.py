#!/usr/bin/env python3
"""Self-checks: tidy-pass behavior plus end-to-end render+verify on every example fixture.

Run: python3 tests/run.py — exits non-zero on the first failed assertion.
"""

import json
import os
import subprocess
import sys
import tempfile

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
sys.path.insert(0, os.path.join(ROOT, "scripts"))
import render as R  # noqa: E402


def load_example(name):
    with open(os.path.join(ROOT, "examples", name)) as f:
        return json.load(f)


def gaps_of(els, ids):
    byid = {e["id"]: e for e in els}
    ms = sorted((byid[i] for i in ids), key=lambda b: b["x"])
    return [b["x"] - (a["x"] + a["w"]) for a, b in zip(ms, ms[1:])]


def test_scr_distribution():
    els = R.normalize_layout(load_example("situation-conflict-resolution.json"))["elements"]
    byid = {e["id"]: e for e in els}
    g = gaps_of(els, ["c1", "c2", "c3"])
    assert abs(g[0] - g[1]) < 1e-6, f"column gaps not equal: {g}"
    tr, c2, c3 = byid["tr"], byid["c2"], byid["c3"]
    mid = ((c2["x"] + c2["w"]) + c3["x"]) / 2
    assert abs((tr["x"] + tr["w"] / 2) - mid) < 1e-6, "triangle not centered in its gap"
    for h, c in (("h1", "c1"), ("h2", "c2"), ("h3", "c3")):
        assert abs((byid[h]["y"] + byid[h]["h"]) - byid[c]["y"]) < 1e-9, f"{h} not flush on {c}"
        assert abs(byid[h]["x"] - byid[c]["x"]) < 1e-9, f"{h} not left-aligned with {c}"
    assert len({round(byid[c]["w"], 6) for c in ("c1", "c2", "c3")}) == 1, "column widths differ"


def test_uneven_untouched():
    layout = {"version": 1, "elements": [
        {"id": "a", "type": "rect", "x": 0.05, "y": 0.40, "w": 0.10, "h": 0.20, "text": "a"},
        {"id": "b", "type": "rect", "x": 0.20, "y": 0.40, "w": 0.10, "h": 0.20, "text": "b"},
        {"id": "c", "type": "rect", "x": 0.80, "y": 0.40, "w": 0.10, "h": 0.20, "text": "c"},
    ]}
    g = gaps_of(R.normalize_layout(layout)["elements"], ["a", "b", "c"])
    assert g[1] / g[0] > 2.5, f"deliberately uneven gaps were equalized: {g}"


def test_bridged_gaps_equalize():
    # arrows between boxes exempt their gaps from the similarity bound
    layout = {"version": 1, "elements": [
        {"id": "a", "type": "rect", "x": 0.05, "y": 0.40, "w": 0.15, "h": 0.20, "text": "a"},
        {"id": "b", "type": "rect", "x": 0.26, "y": 0.40, "w": 0.15, "h": 0.20, "text": "b"},
        {"id": "c", "type": "rect", "x": 0.70, "y": 0.40, "w": 0.15, "h": 0.20, "text": "c"},
        {"id": "ab", "type": "arrow", "from": "a", "to": "b"},
        {"id": "bc", "type": "arrow", "from": "b", "to": "c"},
    ]}
    g = gaps_of(R.normalize_layout(layout)["elements"], ["a", "b", "c"])
    assert abs(g[0] - g[1]) < 1e-6, f"bridged gaps not equalized: {g}"


def _frame(eid, x, w, h=0.6):
    return {"id": eid, "type": "rect", "x": x, "y": 0.2, "w": w, "h": h, "text": eid}


def test_frames_snap_two_thirds():
    layout = {"version": 1, "elements": [_frame("f1", 0.05, 0.55), _frame("f2", 0.68, 0.30)]}
    byid = {e["id"]: e for e in R.normalize_layout(layout)["elements"]}
    ratio = byid["f1"]["w"] / byid["f2"]["w"]
    assert abs(ratio - 2.0) < 1e-9, f"2/3+1/3 not snapped, ratio {ratio}"


def test_frames_quarter_half_quarter():
    layout = {"version": 1, "elements": [
        _frame("f1", 0.05, 0.18), _frame("f2", 0.27, 0.42), _frame("f3", 0.73, 0.20)]}
    byid = {e["id"]: e for e in R.normalize_layout(layout)["elements"]}
    assert abs(byid["f2"]["w"] / byid["f1"]["w"] - 2.0) < 1e-9, "center frame not 2x side frame"
    assert abs(byid["f1"]["w"] - byid["f3"]["w"]) < 1e-9, "side frames not equal"


def test_fifty_fifty_untouched():
    layout = {"version": 1, "elements": [_frame("f1", 0.05, 0.40), _frame("f2", 0.55, 0.40)]}
    byid = {e["id"]: e for e in R.normalize_layout(layout)["elements"]}
    assert abs(byid["f1"]["w"] - byid["f2"]["w"]) < 1e-9, "50/50 frames should stay equal (untouched)"


def test_table_expansion():
    els = R.normalize_layout(load_example("effort-table.json"))["elements"]
    parts = [e for e in els if e.get("id", "").startswith("t1.")]
    kinds = {"rect": 0, "text": 0, "line": 0}
    for e in parts:
        kinds[e["type"]] += 1
    assert kinds == {"rect": 3, "text": 9, "line": 3}, f"unexpected expansion: {kinds}"
    headers = sorted(e["x"] for e in parts if e["type"] == "rect")
    assert abs((headers[1] - headers[0]) - (headers[2] - headers[1])) < 1e-9, "unequal column widths"
    assert not any(e["type"] == "table" for e in els), "table not expanded"


def test_no_native_table_object():
    from pptx import Presentation
    src = os.path.join(ROOT, "examples", "effort-table.json")
    out = os.path.join(tempfile.gettempdir(), "effort-table-check.pptx")
    subprocess.run(
        [sys.executable, os.path.join(ROOT, "scripts", "render.py"), src, "-o", out],
        check=True, capture_output=True,
    )
    prs = Presentation(out)
    assert not any(getattr(sp, "has_table", False) for sp in prs.slides[0].shapes), \
        "native PPT table object found in output"


def test_connector_to_table_rejected():
    layout = {"version": 1, "elements": [
        {"id": "t", "type": "table", "x": 0.1, "y": 0.1, "w": 0.5, "h": 0.4, "rows": [["a"]]},
        {"id": "b", "type": "rect", "x": 0.7, "y": 0.1, "w": 0.2, "h": 0.2, "text": "b"},
        {"id": "arr", "type": "arrow", "from": "t", "to": "b"},
    ]}
    path = os.path.join(tempfile.gettempdir(), "bad-table-conn.json")
    with open(path, "w") as f:
        json.dump(layout, f)
    r = subprocess.run(
        [sys.executable, os.path.join(ROOT, "scripts", "render.py"), path,
         "-o", os.path.join(tempfile.gettempdir(), "bad.pptx")],
        capture_output=True, text=True,
    )
    assert r.returncode == 2 and "cannot attach to tables" in r.stderr, \
        f"connector-to-table not rejected: rc={r.returncode} stderr={r.stderr}"


def test_flush_in_drawn_space():
    # a small sketch upscaled 2.2x: chip drawn 1.2% above its body must still snap
    layout = {"version": 1, "elements": [
        {"id": "chip", "type": "rect", "x": 0.30, "y": 0.300, "w": 0.10, "h": 0.04, "text": "chip"},
        {"id": "body", "type": "rect", "x": 0.30, "y": 0.352, "w": 0.10, "h": 0.15, "bullets": ["…"]},
    ]}
    byid = {e["id"]: e for e in R.normalize_layout(layout)["elements"]}
    gap = byid["body"]["y"] - (byid["chip"]["y"] + byid["chip"]["h"])
    assert abs(gap) < 1e-9, f"chip not flush after upscale: gap {gap}"


def test_flush_survives_size_unification():
    # field test 2026-07-10: chip drawn exactly flush; size unification pools
    # chip height with the content squares and pushes its bottom edge into the
    # frame — flush snap must repair the overlap, not just positive gaps
    layout = {"version": 1, "title": "t", "elements": [
        {"id": "chip", "type": "rect", "x": 0.04, "y": 0.28, "w": 0.18, "h": 0.07, "text": "INPUT", "bold": True},
        {"id": "frame", "type": "rect", "x": 0.04, "y": 0.35, "w": 0.18, "h": 0.50},
        {"id": "s1", "type": "rect", "x": 0.09, "y": 0.40, "w": 0.07, "h": 0.09},
        {"id": "s2", "type": "rect", "x": 0.09, "y": 0.55, "w": 0.07, "h": 0.09},
        {"id": "s3", "type": "rect", "x": 0.09, "y": 0.70, "w": 0.07, "h": 0.09},
    ]}
    byid = {e["id"]: e for e in R.normalize_layout(layout)["elements"]}
    delta = byid["frame"]["y"] - (byid["chip"]["y"] + byid["chip"]["h"])
    assert abs(delta) < 1e-9, f"chip not flush on frame (delta {delta:+.4f})"


def test_header_band_geometry():
    # field test 2026-07-10 (MCK style): header chip becomes a fixed-height band,
    # flush on its frame, never pooled with content-box sizes
    layout = {"version": 1, "title": "t", "elements": [
        {"id": "chip", "type": "rect", "x": 0.04, "y": 0.28, "w": 0.18, "h": 0.07, "text": "INPUT", "header": True},
        {"id": "frame", "type": "rect", "x": 0.04, "y": 0.35, "w": 0.18, "h": 0.50},
        {"id": "s1", "type": "rect", "x": 0.09, "y": 0.40, "w": 0.07, "h": 0.09},
        {"id": "s2", "type": "rect", "x": 0.09, "y": 0.55, "w": 0.07, "h": 0.09},
    ]}
    byid = {e["id"]: e for e in R.normalize_layout(layout)["elements"]}
    assert abs(byid["chip"]["h"] - R.HEADER_H) < 1e-9, f"band height {byid['chip']['h']}"
    delta = byid["frame"]["y"] - (byid["chip"]["y"] + byid["chip"]["h"])
    assert abs(delta) < 1e-9, f"band not flush on frame (delta {delta:+.4f})"


def test_header_band_style():
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    layout = {"version": 1, "elements": [
        {"id": "chip", "type": "rect", "x": 0.1, "y": 0.2, "w": 0.4, "h": 0.08, "text": "INPUT", "header": True},
        {"id": "frame", "type": "rect", "x": 0.1, "y": 0.28, "w": 0.4, "h": 0.5, "text": "body"},
    ]}
    path = os.path.join(tempfile.gettempdir(), "header-style.json")
    with open(path, "w") as f:
        json.dump(layout, f)
    out = os.path.join(tempfile.gettempdir(), "header-style.pptx")
    subprocess.run([sys.executable, os.path.join(ROOT, "scripts", "render.py"), path, "-o", out],
                   check=True, capture_output=True)
    prs = Presentation(out)
    band = next(sp for sp in prs.slides[0].shapes
                if sp.has_text_frame and sp.text_frame.text == "INPUT")
    assert band.fill.fore_color.rgb == RGBColor(0, 0, 0), "band fill not black"
    p = band.text_frame.paragraphs[0]
    assert p.font.color.rgb == RGBColor(255, 255, 255), "band text not white"
    assert p.font.bold, "band text not bold"
    assert p.alignment == PP_ALIGN.LEFT, "band text not left-aligned"


def test_row_wrap_connector_routes_through_gap():
    # last box of row 1 -> first box of row 2: attach bottom->top, not left/right
    boxes = {"r1c": (0.60, 0.20, 0.10, 0.10), "r2a": (0.10, 0.40, 0.10, 0.10)}
    el = {"id": "wrap", "type": "line", "from": "r1c", "to": "r2a"}
    p1, p2, kind, s1, s2 = R.resolve_connector(el, boxes)
    assert (s1, s2) == (R.SIDE_BOTTOM, R.SIDE_TOP), f"sides {s1},{s2}"
    assert abs(p1[0] - 0.65) < 1e-9 and abs(p1[1] - 0.30) < 1e-9, f"p1 {p1}"
    assert abs(p2[0] - 0.15) < 1e-9 and abs(p2[1] - 0.40) < 1e-9, f"p2 {p2}"
    assert kind == "elbow"
    # same-row neighbors keep left/right attachment
    boxes = {"a": (0.10, 0.20, 0.10, 0.10), "b": (0.40, 0.21, 0.10, 0.10)}
    _, _, _, s1, s2 = R.resolve_connector({"type": "line", "from": "a", "to": "b"}, boxes)
    assert (s1, s2) == (R.SIDE_RIGHT, R.SIDE_LEFT), f"same-row sides {s1},{s2}"


def test_header_seats_from_inside_frame():
    # field test 2026-07-10 (second run): band transcribed as the frame's top
    # compartment (band top == frame top) must seat flush ON the frame top
    layout = {"version": 1, "title": "t", "elements": [
        {"id": "band", "type": "rect", "x": 0.06, "y": 0.30, "w": 0.20, "h": 0.09, "text": "INPUT", "header": True},
        {"id": "frame", "type": "rect", "x": 0.06, "y": 0.30, "w": 0.20, "h": 0.55},
        {"id": "s1", "type": "rect", "x": 0.10, "y": 0.45, "w": 0.08, "h": 0.10},
    ]}
    byid = {e["id"]: e for e in R.normalize_layout(layout)["elements"]}
    delta = byid["frame"]["y"] - (byid["band"]["y"] + byid["band"]["h"])
    assert abs(delta) < 1e-9, f"band not seated on frame top (delta {delta:+.4f})"
    assert abs(byid["band"]["h"] - R.HEADER_H) < 1e-9


def test_headers_render_on_top():
    from pptx import Presentation
    # band listed first and overlapping the frame: z-order must still put it above
    layout = {"version": 1, "elements": [
        {"id": "band", "type": "rect", "x": 0.10, "y": 0.20, "w": 0.40, "h": 0.08, "text": "SYS", "header": True},
        {"id": "frame", "type": "rect", "x": 0.10, "y": 0.20, "w": 0.40, "h": 0.60},
    ]}
    path = os.path.join(tempfile.gettempdir(), "zorder.json")
    with open(path, "w") as f:
        json.dump(layout, f)
    out = os.path.join(tempfile.gettempdir(), "zorder.pptx")
    subprocess.run([sys.executable, os.path.join(ROOT, "scripts", "render.py"), path, "-o", out],
                   check=True, capture_output=True)
    order = [sp.text_frame.text if sp.has_text_frame else "" for sp in Presentation(out).slides[0].shapes]
    assert order.index("SYS") > order.index(""), f"band not on top: {order}"


def test_frames_share_top_and_bottom():
    # includes a stacked column of small boxes inside f1: its tall UNION must
    # not count as a frame (that would make 4 candidates and disable the pass)
    layout = {"version": 1, "elements": [
        {"id": "f1", "type": "rect", "x": 0.05, "y": 0.28, "w": 0.24, "h": 0.56, "text": "a"},
        {"id": "s1", "type": "rect", "x": 0.10, "y": 0.34, "w": 0.06, "h": 0.08},
        {"id": "s2", "type": "rect", "x": 0.10, "y": 0.52, "w": 0.06, "h": 0.08},
        {"id": "s3", "type": "rect", "x": 0.10, "y": 0.70, "w": 0.06, "h": 0.08},
        {"id": "f2", "type": "rect", "x": 0.38, "y": 0.30, "w": 0.24, "h": 0.52, "text": "b"},
        {"id": "f3", "type": "rect", "x": 0.71, "y": 0.32, "w": 0.24, "h": 0.48, "text": "c"},
    ]}
    byid = {e["id"]: e for e in R.normalize_layout(layout)["elements"]}
    tops = {round(byid[f]["y"], 9) for f in ("f1", "f2", "f3")}
    hts = {round(byid[f]["h"], 9) for f in ("f1", "f2", "f3")}
    assert len(tops) == 1, f"frame tops differ: {tops}"
    assert len(hts) == 1, f"frame heights differ: {hts}"


def test_end_to_end():
    for name in sorted(os.listdir(os.path.join(ROOT, "examples"))):
        if not name.endswith(".json"):
            continue
        src = os.path.join(ROOT, "examples", name)
        out = os.path.join(tempfile.gettempdir(), name.replace(".json", ".pptx"))
        subprocess.run(
            [sys.executable, os.path.join(ROOT, "scripts", "render.py"), src, "-o", out],
            check=True, capture_output=True,
        )
        r = subprocess.run(
            [sys.executable, os.path.join(ROOT, "scripts", "verify.py"), src, out],
            capture_output=True, text=True,
        )
        v = json.loads(r.stdout)
        assert v["passed"], f"{name}: verify failed: {v['failures']}"


if __name__ == "__main__":
    for check in (test_scr_distribution, test_uneven_untouched,
                  test_bridged_gaps_equalize, test_flush_in_drawn_space,
                  test_flush_survives_size_unification,
                  test_header_band_geometry, test_header_band_style,
                  test_header_seats_from_inside_frame, test_headers_render_on_top,
                  test_frames_share_top_and_bottom,
                  test_row_wrap_connector_routes_through_gap,
                  test_frames_snap_two_thirds, test_frames_quarter_half_quarter,
                  test_fifty_fifty_untouched,
                  test_table_expansion, test_no_native_table_object,
                  test_connector_to_table_rejected, test_end_to_end):
        check()
        print(f"ok {check.__name__}")
    print("all checks passed")

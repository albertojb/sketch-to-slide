#!/usr/bin/env python3
"""sketch-to-slide verifier: checks a rendered .pptx against its layout JSON.

Machine-readable gate — prints JSON with a `passed` boolean; exit code 0 only when passed.
"""

import argparse
import json
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import render as R  # noqa: E402

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.shapes.connector import Connector

TOL = 0.03


def near(a, b, tol=TOL):
    return abs(a[0] - b[0]) <= tol and abs(a[1] - b[1]) <= tol


def rendered_inventory(path):
    prs = Presentation(path)
    sw, sh = prs.slide_width, prs.slide_height
    slide = prs.slides[0]
    shapes, conns = [], []
    for sp in slide.shapes:
        if isinstance(sp, Connector):
            ln = sp.line._get_or_add_ln()
            conns.append({
                "p1": (sp.begin_x / sw, sp.begin_y / sh),
                "p2": (sp.end_x / sw, sp.end_y / sh),
                "head": ln.find(qn("a:headEnd")) is not None,
                "tail": ln.find(qn("a:tailEnd")) is not None,
                "used": False,
            })
        else:
            try:
                auto = sp.auto_shape_type
            except Exception:
                auto = None
            shapes.append({
                "auto": auto,
                "textbox": sp.shape_type == MSO_SHAPE_TYPE.TEXT_BOX,
                "center": ((sp.left + sp.width / 2) / sw, (sp.top + sp.height / 2) / sh),
                "text": sp.text_frame.text.strip() if sp.has_text_frame else "",
                "used": False,
            })
    return shapes, conns


def expected_items(layout):
    boxes = R.boxes_by_id(layout)
    exp_shapes, exp_conns = [], []
    if layout.get("title"):
        x, y, w, h = R.TITLE_BOX
        exp_shapes.append({"what": "title", "auto": None, "textbox": True,
                           "center": (x + w / 2, y + h / 2), "text": str(layout["title"]).strip()})
    for el in layout["elements"]:
        t = el["type"]
        name = el.get("id") or t
        if t in R.SHAPE_MAP:
            exp_shapes.append({"what": name, "auto": R.SHAPE_MAP[t], "textbox": False,
                               "center": (el["x"] + el["w"] / 2, el["y"] + el["h"] / 2),
                               "text": str(el.get("text", "")).strip()})
        elif t == "text":
            exp_shapes.append({"what": name, "auto": None, "textbox": True,
                               "center": (el["x"] + el["w"] / 2, el["y"] + el["h"] / 2),
                               "text": str(el.get("text", "")).strip()})
        else:
            p1, p2 = R.resolve_connector(el, boxes)
            exp_conns.append({"what": name, "p1": p1, "p2": p2,
                              "head": t == "double_arrow",
                              "tail": t in ("arrow", "double_arrow")})
            if el.get("text"):
                mid = ((p1[0] + p2[0]) / 2, (p1[1] + p2[1]) / 2)
                exp_shapes.append({"what": f"label:{name}", "auto": None, "textbox": True,
                                   "center": mid, "text": str(el["text"]).strip()})
    return exp_shapes, exp_conns


def match(layout_path, pptx_path):
    layout = R.load_layout(layout_path)
    exp_shapes, exp_conns = expected_items(layout)
    got_shapes, got_conns = rendered_inventory(pptx_path)
    failures = []

    for e in exp_shapes:
        hit = next((g for g in got_shapes
                    if not g["used"] and g["auto"] == e["auto"] and g["textbox"] == e["textbox"]
                    and near(g["center"], e["center"]) and g["text"] == e["text"]), None)
        if hit:
            hit["used"] = True
        else:
            failures.append(f"shape '{e['what']}' not found (type, position, or text mismatch)")

    for e in exp_conns:
        hit = next((g for g in got_conns
                    if not g["used"] and g["head"] == e["head"] and g["tail"] == e["tail"]
                    and near(g["p1"], e["p1"]) and near(g["p2"], e["p2"])), None)
        if hit:
            hit["used"] = True
        else:
            failures.append(f"connector '{e['what']}' not found (endpoints or arrowheads mismatch)")

    extra_shapes = sum(1 for g in got_shapes if not g["used"])
    extra_conns = sum(1 for g in got_conns if not g["used"])
    if extra_shapes:
        failures.append(f"{extra_shapes} rendered shape(s) not present in layout")
    if extra_conns:
        failures.append(f"{extra_conns} rendered connector(s) not present in layout")

    return {
        "passed": not failures,
        "expected": {"shapes": len(exp_shapes), "connectors": len(exp_conns)},
        "rendered": {"shapes": len(got_shapes), "connectors": len(got_conns)},
        "failures": failures,
    }


def main():
    ap = argparse.ArgumentParser(description="Verify a rendered .pptx matches its layout JSON")
    ap.add_argument("layout")
    ap.add_argument("pptx")
    args = ap.parse_args()
    result = match(args.layout, args.pptx)
    print(json.dumps(result, indent=2))
    sys.exit(0 if result["passed"] else 1)


if __name__ == "__main__":
    main()
